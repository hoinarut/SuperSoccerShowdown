#!/usr/bin/env python3
"""Generate SSS architecture presentation (~25 min)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

OUT = Path(__file__).resolve().parents[1] / "docs" / "SSS_Architecture_Presentation.pptx"

# Palette — deep teal / slate (not purple/cream AI defaults)
NAVY = RGBColor(0x0F, 0x2C, 0x3A)
TEAL = RGBColor(0x1A, 0x6B, 0x6B)
ACCENT = RGBColor(0xE0, 0x7A, 0x3D)
SLATE = RGBColor(0x3D, 0x4F, 0x5C)
LIGHT = RGBColor(0xF5, 0xF7, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6B, 0x7C, 0x88)
CARD_BG = RGBColor(0xE8, 0xEE, 0xF0)


def set_run(run, text, size=18, bold=False, color=SLATE, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_bg(slide, color=LIGHT):
    fill = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    fill.fill.solid()
    fill.fill.fore_color.rgb = color
    fill.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = fill._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_accent_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.12), Inches(7.5)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def add_footer(slide, page, total=14):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(11.5), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    set_run(run, f"SSS · Architecture walkthrough  ·  {page}/{total}", size=11, color=MUTED)
    p2 = tf.add_paragraph()
    # page number right via separate box
    num = slide.shapes.add_textbox(Inches(12.2), Inches(7.05), Inches(0.8), Inches(0.35))
    np = num.text_frame.paragraphs[0]
    np.alignment = PP_ALIGN.RIGHT
    r = np.add_run()
    set_run(r, str(page), size=11, color=MUTED)


def title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.9), Inches(13.333), Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, "SSS", size=48, bold=True, color=WHITE, font="Calibri")

    sub = slide.shapes.add_textbox(Inches(0.9), Inches(3.2), Inches(11), Inches(1))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, "Domain, Architecture & Testability", size=28, color=RGBColor(0xA8, 0xC5, 0xC5))

    meta = slide.shapes.add_textbox(Inches(0.9), Inches(6.2), Inches(11), Inches(0.8))
    tf = meta.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, "Fantasy team builder  ·  .NET 10 hexagonal backend  ·  ~25 minutes", size=16, color=MUTED)


def section_slide(prs, number, title, subtitle, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_footer(slide, page)

    num = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11), Inches(0.5))
    p = num.text_frame.paragraphs[0]
    r = p.add_run()
    set_run(r, number, size=18, bold=True, color=ACCENT)

    t = slide.shapes.add_textbox(Inches(0.9), Inches(2.9), Inches(11), Inches(1))
    p = t.text_frame.paragraphs[0]
    r = p.add_run()
    set_run(r, title, size=36, bold=True, color=WHITE)

    s = slide.shapes.add_textbox(Inches(0.9), Inches(4.0), Inches(11), Inches(0.6))
    p = s.text_frame.paragraphs[0]
    r = p.add_run()
    set_run(r, subtitle, size=18, color=RGBColor(0xA8, 0xC5, 0xC5))


def add_title(slide, text):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    set_run(r, text, size=28, bold=True, color=NAVY)


def add_bullets(slide, left, top, width, height, items, size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (level, text) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.space_after = Pt(6)
        r = p.add_run()
        set_run(r, text, size=size - level * 2, color=SLATE if level == 0 else MUTED)


def card(slide, left, top, width, height, title, body_lines, title_color=TEAL):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = CARD_BG
    shape.adjustments[0] = 0.08

    tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.15), Inches(width - 0.4), Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    set_run(r, title, size=14, bold=True, color=title_color)

    body = slide.shapes.add_textbox(
        Inches(left + 0.2), Inches(top + 0.55), Inches(width - 0.4), Inches(height - 0.7)
    )
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(body_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        r = p.add_run()
        set_run(r, line, size=13, color=SLATE)


def overview_slide(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_footer(slide, page)
    add_title(slide, "What SSS does")

    add_bullets(
        slide,
        0.5,
        1.1,
        7.5,
        3.5,
        [
            (0, "Fantasy sports team builder over external character APIs"),
            (0, "Pick a Universe (Star Wars or Pokémon), name a team, set attackers/defenders"),
            (0, "System builds a 5-player roster and assigns Goalie / Defence / Offence"),
            (0, "Stack: .NET 10 Web API · Mediator CQRS · EF Core · Angular · AWS Lambda"),
        ],
        size=17,
    )

    card(
        slide,
        8.5,
        1.2,
        4.2,
        4.5,
        "Talk agenda (~25 min)",
        [
            "1. Domain & data structures",
            "2. Architecture & patterns",
            "3. Tests & hexagonal layout",
            "",
            "Goal: leave knowing where",
            "rules live, how ports",
            "isolate adapters, and how",
            "each layer is tested.",
        ],
        title_color=ACCENT,
    )


def domain_model_slide(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_footer(slide, page)
    add_title(slide, "Domain model")

    # Relationship diagram as cards
    card(slide, 0.5, 1.15, 3.8, 2.4, "Universe", [
        "Reference data",
        "Name, IsEnabled, ApiUrl",
        "Seeded: Pokémon, Star Wars",
        "",
        "1 ──< many Teams",
    ])
    card(slide, 4.7, 1.15, 3.8, 2.4, "Team  (aggregate root)", [
        "Name, UniverseId",
        "AttackersCount / DefendersCount",
        "List<Player> roster (max 5)",
        "",
        "Owns add / assign / validate",
    ], title_color=ACCENT)
    card(slide, 8.9, 1.15, 3.8, 2.4, "Player", [
        "Name, Height, Weight",
        "PlayerType enum",
        "ExternalResourceId",
        "",
        "Child of Team",
    ])

    add_bullets(
        slide,
        0.5,
        3.8,
        12,
        2.8,
        [
            (0, "BaseEntity → int Id; no separate value objects (height/weight are double)"),
            (0, "PlayerType: Goalie | Defence | Offence"),
            (0, "Constants: NumberOfPlayers = 5; shared error messages"),
            (0, "Data structures: List<Player>, List<int> resource IDs, IEnumerable<IUniverseDataService>, IQueryable via repos"),
        ],
        size=16,
    )


def domain_rules_slide(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_footer(slide, page)
    add_title(slide, "Domain rules live on Team")

    card(slide, 0.5, 1.15, 4.0, 5.2, "AddPlayer", [
        "Guards roster size",
        "Throws when count == 5",
        "",
        "Encapsulation:",
        "private setters + behavior",
        "methods — not anemic",
    ])
    card(slide, 4.7, 1.15, 4.0, 5.2, "SetPlayerTypes", [
        "1. Tallest → Goalie",
        "   (tie: name ascending)",
        "2. Heaviest × Defenders",
        "   → Defence",
        "3. Shortest × Attackers",
        "   → Offence",
        "",
        "Uses UnassignedPlayers",
        "(Type == 0 filter)",
    ], title_color=ACCENT)
    card(slide, 8.9, 1.15, 4.0, 5.2, "ValidateSetup", [
        "Exactly 5 players",
        "Exactly 1 Goalie",
        "Defence count matches",
        "  DefendersCount",
        "Offence count matches",
        "  AttackersCount",
        "",
        "Invariant check before",
        "persistence",
    ])


def hex_arch_slide(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_footer(slide, page)
    add_title(slide, "Hexagonal architecture")

    # Layer boxes
    layers = [
        (0.5, "Api", "Driving adapter", "Controllers · FluentValidation\nMediator send · composition root"),
        (3.7, "Application", "Use cases", "Handlers (CQRS)\nTeamManager orchestration"),
        (6.9, "Domain", "Core", "Entities · enums\nPorts (interfaces)"),
        (10.1, "Infrastructure", "Driven adapters", "EF repos · HTTP clients\nSWAPI / PokéAPI"),
    ]
    for left, name, role, body in layers:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.2), Inches(2.9), Inches(3.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE if name != "Domain" else NAVY
        shape.line.color.rgb = TEAL if name == "Domain" else CARD_BG
        shape.adjustments[0] = 0.08

        tc = WHITE if name == "Domain" else NAVY
        rc = ACCENT if name == "Domain" else TEAL
        bc = RGBColor(0xC5, 0xD5, 0xD5) if name == "Domain" else SLATE

        t = slide.shapes.add_textbox(Inches(left + 0.15), Inches(1.4), Inches(2.6), Inches(0.4))
        p = t.text_frame.paragraphs[0]
        r = p.add_run()
        set_run(r, name, size=16, bold=True, color=tc)

        role_box = slide.shapes.add_textbox(Inches(left + 0.15), Inches(1.85), Inches(2.6), Inches(0.35))
        p = role_box.text_frame.paragraphs[0]
        r = p.add_run()
        set_run(r, role, size=12, bold=True, color=rc)

        b = slide.shapes.add_textbox(Inches(left + 0.15), Inches(2.4), Inches(2.6), Inches(2.2))
        tf = b.text_frame
        tf.word_wrap = True
        for i, line in enumerate(body.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(4)
            r = p.add_run()
            set_run(r, line, size=12, color=bc)

    note = slide.shapes.add_textbox(Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.5))
    tf = note.text_frame
    tf.word_wrap = True
    lines = [
        "Dependency direction:  Api → Application → Domain ← Infrastructure",
        "Ports live in Domain (IRepository, IUniverseDataService, ITeamManager). Adapters implement them in Infrastructure.",
        "Domain has zero project/NuGet dependencies — the innermost hexagon.",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        r = p.add_run()
        set_run(r, line, size=14, color=SLATE)


def flow_slide(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_footer(slide, page)
    add_title(slide, "Create-team flow")

    steps = [
        ("1", "HTTP", "POST /teams\nFluentValidation"),
        ("2", "Mediator", "CreateTeamCommand\n→ Handler"),
        ("3", "TeamManager", "Pick strategy\nFetch & build roster"),
        ("4", "Ports", "IUniverseDataService\nI*Repository"),
        ("5", "Persist", "EF Core save\nreturn TeamDto"),
    ]
    for i, (num, title, body) in enumerate(steps):
        left = 0.4 + i * 2.55
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.5), Inches(2.3), Inches(2.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = CARD_BG
        shape.adjustments[0] = 0.1

        n = slide.shapes.add_textbox(Inches(left + 0.15), Inches(1.65), Inches(2.0), Inches(0.4))
        p = n.text_frame.paragraphs[0]
        r = p.add_run()
        set_run(r, num, size=22, bold=True, color=ACCENT)

        t = slide.shapes.add_textbox(Inches(left + 0.15), Inches(2.2), Inches(2.0), Inches(0.4))
        p = t.text_frame.paragraphs[0]
        r = p.add_run()
        set_run(r, title, size=15, bold=True, color=NAVY)

        b = slide.shapes.add_textbox(Inches(left + 0.15), Inches(2.7), Inches(2.0), Inches(1.4))
        tf = b.text_frame
        tf.word_wrap = True
        for j, line in enumerate(body.split("\n")):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            r = p.add_run()
            set_run(r, line, size=13, color=SLATE)

        if i < len(steps) - 1:
            arrow = slide.shapes.add_textbox(Inches(left + 2.15), Inches(2.6), Inches(0.4), Inches(0.4))
            p = arrow.text_frame.paragraphs[0]
            r = p.add_run()
            set_run(r, "→", size=20, bold=True, color=TEAL)

    add_bullets(
        slide,
        0.5,
        4.7,
        12,
        1.8,
        [
            (0, "Strategy selection: IEnumerable<IUniverseDataService>.FirstOrDefault(s => s.CanHandle(universe))"),
            (0, "Random unused external IDs → GeneratePlayer → AddPlayer → SetPlayerTypes → ValidateSetup → save"),
            (0, "New universe = new IUniverseDataService + DI registration — TeamManager stays closed for modification"),
        ],
        size=15,
    )


def solid_slide(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_footer(slide, page)
    add_title(slide, "SOLID, DRY & patterns")

    items = [
        ("S", "Single Responsibility", "Controllers = HTTP; Team = roster rules; handlers = mapping"),
        ("O", "Open/Closed", "New universe via IUniverseDataService — no TeamManager change"),
        ("L", "Liskov", "SW / Pokémon services & repos honor their ports interchangeably"),
        ("I", "Interface Segregation", "Narrow ports: Exists, GetUsedExternalResourceIds — not one fat API"),
        ("D", "Dependency Inversion", "Application depends on Domain ports; Infra supplies adapters"),
    ]
    for i, (letter, name, desc) in enumerate(items):
        top = 1.05 + i * 0.85
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(top), Inches(0.7), Inches(0.7)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = TEAL
        badge.line.fill.background()
        badge.adjustments[0] = 0.15
        lt = slide.shapes.add_textbox(Inches(0.5), Inches(top + 0.12), Inches(0.7), Inches(0.5))
        p = lt.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        set_run(r, letter, size=20, bold=True, color=WHITE)

        t = slide.shapes.add_textbox(Inches(1.4), Inches(top + 0.05), Inches(5), Inches(0.35))
        p = t.text_frame.paragraphs[0]
        r = p.add_run()
        set_run(r, name, size=15, bold=True, color=NAVY)

        d = slide.shapes.add_textbox(Inches(1.4), Inches(top + 0.35), Inches(5.5), Inches(0.4))
        p = d.text_frame.paragraphs[0]
        r = p.add_run()
        set_run(r, desc, size=13, color=SLATE)

    card(
        slide,
        7.5,
        1.15,
        5.3,
        5.3,
        "Named patterns in the codebase",
        [
            "• Ports & Adapters (hexagonal)",
            "• Repository + RepositoryBase<T>",
            "• Strategy (CanHandle universe)",
            "• CQRS + Mediator handlers",
            "• Aggregate (Team boundary)",
            "• Anti-corruption DTOs",
            "  (StarWarsPerson → Player)",
            "• Composition Root (Program.cs)",
            "",
            "DRY: shared Constants, seed data,",
            "RepositoryBase, test base classes",
        ],
        title_color=ACCENT,
    )


def tests_slide(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_footer(slide, page)
    add_title(slide, "What is tested — and how")

    card(slide, 0.5, 1.15, 6.0, 2.5, "Domain  ·  pure unit", [
        "TeamTests — roster limits, role algorithm,",
        "tie-breaks, ValidateSetup (~14 cases)",
        "PlayerTests — measurements & type",
        "No DB, no HTTP, no mocks needed",
    ])
    card(slide, 6.8, 1.15, 6.0, 2.5, "Application  ·  port mocks", [
        "TeamManagerTests (~10) with Moq",
        "Happy path, exhausted resources,",
        "skip invalid/dup names, strategy pick",
        "Strict mocks via BaseUnitTests",
    ], title_color=ACCENT)

    card(slide, 0.5, 3.9, 6.0, 2.5, "Infrastructure  ·  adapter isolation", [
        "EF repos: InMemory DB integration",
        "Universe HTTP: mocked HttpMessageHandler",
        "Mapping, CanHandle, error paths",
    ])
    card(slide, 6.8, 3.9, 6.0, 2.5, "Api  ·  pipeline integration", [
        "WebApplicationFactory + InMemory DB",
        "ExternalApiMockHandler for SW/Pokémon",
        "Real controllers → Mediator → DI",
        "Stack: xUnit · FluentAssertions · Moq",
    ], title_color=TEAL)


def hex_testability_slide(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_footer(slide, page)
    add_title(slide, "Why hexagonal helps testability")

    add_bullets(
        slide,
        0.5,
        1.2,
        12.5,
        5.5,
        [
            (0, "Domain has no dependencies → fastest feedback on business rules"),
            (0, "Application depends only on ports → replace DB/HTTP with Moq in TeamManagerTests"),
            (0, "Infrastructure adapters tested alone → fake HTTP handler / InMemory EF"),
            (0, "Api factory swaps outer rings without rewriting use cases"),
            (0, "Same ports that enable OCP also enable the test pyramid (~51 tests across layers)"),
            (0, ""),
            (0, "Structure to remember:"),
            (1, "tests/MyApp.Domain.Tests        → entities (pure)"),
            (1, "tests/MyApp.Application.Tests   → services with mocked ports"),
            (1, "tests/MyApp.Infrastructure.Tests → EF + external API adapters"),
            (1, "tests/MyApp.Api.Tests            → full HTTP pipeline"),
        ],
        size=16,
    )


def recap_slide(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_accent_bar(slide)
    add_footer(slide, page)
    add_title(slide, "Recap")

    points = [
        ("Domain", "Universe–Team–Player; Team aggregate owns roster invariants and role assignment"),
        ("Architecture", "Ports in Domain, adapters in Infrastructure; Api drives via Mediator CQRS"),
        ("Principles", "SOLID via Strategy + DIP; DRY via shared bases/constants; OCP for new universes"),
        ("Tests", "Layer-aligned pyramid; hexagonal boundaries make each ring independently testable"),
    ]
    for i, (label, text) in enumerate(points):
        top = 1.2 + i * 1.25
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(top), Inches(12.3), Inches(1.1)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = CARD_BG
        shape.adjustments[0] = 0.08

        lb = slide.shapes.add_textbox(Inches(0.75), Inches(top + 0.2), Inches(2.2), Inches(0.7))
        p = lb.text_frame.paragraphs[0]
        r = p.add_run()
        set_run(r, label, size=16, bold=True, color=TEAL)

        tx = slide.shapes.add_textbox(Inches(3.0), Inches(top + 0.25), Inches(9.4), Inches(0.7))
        p = tx.text_frame.paragraphs[0]
        r = p.add_run()
        set_run(r, text, size=15, color=SLATE)


def qa_slide(prs, page):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)

    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(11), Inches(1))
    p = box.text_frame.paragraphs[0]
    r = p.add_run()
    set_run(r, "Questions?", size=44, bold=True, color=WHITE)

    sub = slide.shapes.add_textbox(Inches(0.9), Inches(3.8), Inches(11), Inches(1.2))
    tf = sub.text_frame
    tf.word_wrap = True
    lines = [
        "Handy paths: Domain/Entities/Team.cs  ·  Domain/Ports/  ·  Application/Services/TeamManager.cs",
        "Adapters: Infrastructure/ExternalServices/*  ·  Tests mirror src layers under tests/",
    ]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        r = p.add_run()
        set_run(r, line, size=15, color=RGBColor(0xA8, 0xC5, 0xC5))

    add_footer(slide, page)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)                          # 1
    overview_slide(prs, 2)                    # 2
    section_slide(prs, "01", "Domain & data structures",
                  "How the fantasy-team domain is modelled", 3)
    domain_model_slide(prs, 4)                # 4
    domain_rules_slide(prs, 5)                # 5
    section_slide(prs, "02", "Architecture & design patterns",
                  "Hexagonal layout, SOLID, and the patterns that show up in code", 6)
    hex_arch_slide(prs, 7)                    # 7
    flow_slide(prs, 8)                        # 8
    solid_slide(prs, 9)                       # 9
    section_slide(prs, "03", "Tests & testability",
                  "What each layer proves — and how ports unlock the pyramid", 10)
    tests_slide(prs, 11)                      # 11
    hex_testability_slide(prs, 12)            # 12
    recap_slide(prs, 13)                      # 13
    qa_slide(prs, 14)                         # 14

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
